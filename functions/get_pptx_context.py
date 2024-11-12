import os
from io import BytesIO
from typing import Annotated

import openai
import pytesseract
from livekit.agents import llm
from PIL import Image, UnidentifiedImageError
from pptx import Presentation


class GetPPTXContext(llm.FunctionContext):
    """
    Function to answer questions based on content from all .pptx files in assets/powerpoint.
    """

    @llm.ai_callable()
    async def answer_pptx_question(
        self,
        question: Annotated[
            str, llm.TypeInfo(description="The question to search for in the .pptx files")
        ],
    ):
        """Answers questions based on content from all .pptx files in assets/powerpoint, including text and images with OCR."""
        
        pptx_directory = "assets/powerpoint"
        if not os.path.exists(pptx_directory):
            raise FileNotFoundError(f"The directory {pptx_directory} was not found.")
        
        extracted_content = []

        # Loop through all .pptx files in the directory
        for pptx_file in os.listdir(pptx_directory):
            if pptx_file.endswith(".pptx"):
                pptx_path = os.path.join(pptx_directory, pptx_file)
                prs = Presentation(pptx_path)
                
                for slide_num, slide in enumerate(prs.slides, start=1):
                    slide_content = {"file_name": pptx_file, "slide_number": slide_num, "text": "", "images": []}
                    
                    # Extract text
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slide_content["text"] += shape.text + " "
                    
                        # Extract images and perform OCR
                        if shape.shape_type == 13:  # Shape with a picture
                            try:
                                image = shape.image
                                img_bytes = BytesIO(image.blob)
                                img = Image.open(img_bytes)
                                
                                # OCR to extract text from image
                                image_text = pytesseract.image_to_string(img)
                                slide_content["images"].append({"text": image_text})
                            
                            except UnidentifiedImageError:
                                print(f"Skipping unsupported image format on slide {slide_num} in file {pptx_file}")
                            except TypeError as e:
                                print(f"Skipping image due to error: {e} on slide {slide_num} in file {pptx_file}")
                    
                    extracted_content.append(slide_content)
        
        # Combine extracted text for context from all files
        context_text = ""
        for slide in extracted_content:
            context_text += f"File: {slide['file_name']} - Slide {slide['slide_number']}:\n{slide['text']}\n"
            for image_info in slide['images']:
                context_text += f"Image text (from OCR): {image_info['text']}\n"
        
        # Use OpenAI to answer the question based on the combined context
        openai.api_key = os.getenv("OPENAI_API_KEY")
        response = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are an assistant that answers questions based on the provided context."},
                {"role": "user", "content": f"Context: {context_text}"},
                {"role": "user", "content": f"Question: {question}"}
            ],
            max_tokens=150,
            temperature=0
        )
        
        answer = response.choices[0].message.content.strip()
        return answer